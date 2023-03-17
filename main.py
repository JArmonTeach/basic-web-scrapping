from bs4 import BeautifulSoup
import requests, os
from pptx import Presentation
from pptx.util import Inches, Pt

# Send a GET request to the website and parse the HTML content
url = input("Enter the URL of the USCCB website with the readings:")
html_text = requests.get(url)
soup = BeautifulSoup(html_text.content, 'html.parser')

# Locate the HTML elements that contain the text you want to scrape
text = ''
for element in soup.find_all(class_='content-body'):
    text += str(element.prettify())

# Replace all <br> tags with newline characters
soup = BeautifulSoup(text, 'html.parser')
for br in soup.find_all('br'):
    br.replace_with('\n')

# Split the text into paragraphs based on newline characters
lines = str(soup).split('\n')
paragraphs = []
paragraph = ''
for line in lines:
    if line:
        # Check if the line contains any of the phrases you want to remove; you can add or change the text that needs to be filtered
        if '<div class="content-body">' in line or '</div>' in line or 'R. (1)' in line or 'R.' in line or '<strong>' in line or '</strong>' in line:
            continue
        paragraph += line
    else:
        paragraphs.append(paragraph.strip())
        paragraph = ''
if paragraph:
    paragraphs.append(paragraph.strip())

# Create a new PowerPoint presentation, set the size of the slide and add the first slide
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Create a text box on the slide
left = Inches(1)
top = Inches(.55)
width = Inches(14)
height = Inches(7.98)
textbox = slide.shapes.add_textbox(left, top, width, height)

# Reformat the text
font_name = 'Georgia'
font_size = Pt(36)
textbox.text = paragraphs[0]
font = textbox.text_frame.paragraphs[0].font
font.name = font_name
font.size = font_size

# Add the scraped text to slide(s)
line_count = textbox.text_frame.paragraphs[0].text.count('\n') + 1
for paragraph in paragraphs[1:]:
    # Check if the text in the current slide has exceeded 13 lines
    if line_count >= 13:
        # Create a new slide and text box
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        font = textbox.text_frame.paragraphs[0].font
        font.name = font_name
        font.size = font_size
        line_count = 0
    # Add the paragraph to the current slide's text box
    p = textbox.text_frame.add_paragraph()
    p.text = paragraph
    font = p.font
    font.name = font_name
    font.size = font_size
    line_count += paragraph.count('\n') + 1

# Save the PowerPoint presentation in the Downloads folder with the filename 'Day_Readings.pptx'
downloads_path = os.path.join(os.path.expanduser('~'), 'Downloads')
prs.save(os.path.join(downloads_path, 'Day_Readings.pptx'))
